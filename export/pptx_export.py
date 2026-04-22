"""
PowerPoint (PPTX) executive deck export for SOC-Report.

Generates a condensed 8-10 slide presentation suitable for management briefings.
Uses python-pptx; no PowerPoint template file required (fully programmatic).

Slide structure:
  1. Title slide (customer + Logicalis logos, report period)
  2. Executive Summary (key metric callouts)
  3. Incident Severity (chart)
  4. Incident Trend (12-month chart)
  5. Incident Status / Resolution (chart)
  6. Top Alerts (chart + top-5 table)
  7. Pending Tickets (table, if data present)
  8. SOCRadar Threat Intelligence (if SOCRadar data present)
  9. Recommendations (table)
 10. Monitoring Scope + Confidentiality
"""

import os
import re
import logging
from io import BytesIO
from datetime import datetime

import markdown as md_lib
from bs4 import BeautifulSoup

logger = logging.getLogger(__name__)

# Logicalis brand colours (RGB tuples)
_BLUE  = (0x1F, 0x6F, 0xEB)
_RED   = (0xDC, 0x26, 0x26)   # Logicalis red (used for title bars and accents)
_DARK  = (0x1A, 0x1A, 0x2E)
_WHITE = (0xFF, 0xFF, 0xFF)
_GREY  = (0x64, 0x74, 0x8B)
_LIGHT_BG = (0xF8, 0xFA, 0xFC)


def _rgb(r, g, b):
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)


def _pt(n):
    from pptx.util import Pt
    return Pt(n)


def _inches(n):
    from pptx.util import Inches
    return Inches(n)


def _get_report_month(report_date: str) -> str:
    try:
        parts = report_date.split(" to ")
        dt = datetime.strptime(parts[-1].strip(), "%Y-%m-%d")
        return dt.strftime("%B %Y")
    except Exception:
        return report_date


def _extract_section_text(soup: BeautifulSoup, heading_keyword: str,
                          max_chars: int = 600) -> str:
    """
    Extract plain text for a section identified by a keyword in its heading.
    Returns up to max_chars characters.
    """
    target = None
    for tag in soup.find_all(re.compile(r"^h[1-4]$")):
        if heading_keyword.lower() in tag.get_text().lower():
            target = tag
            break
    if not target:
        return ""

    parts = []
    for sibling in target.find_next_siblings():
        if re.match(r"^h[1-4]$", sibling.name or ""):
            break
        text = sibling.get_text(separator=" ").strip()
        if text:
            parts.append(text)
        if sum(len(p) for p in parts) >= max_chars:
            break

    full = " ".join(parts)
    return full[:max_chars] + ("…" if len(full) > max_chars else "")


def _extract_table_rows(soup: BeautifulSoup, heading_keyword: str,
                        max_rows: int = 6) -> tuple[list[str], list[list[str]]]:
    """
    Find the first table after a heading matching heading_keyword.
    Returns (headers, rows).
    """
    target = None
    for tag in soup.find_all(re.compile(r"^h[1-4]$")):
        if heading_keyword.lower() in tag.get_text().lower():
            target = tag
            break
    if not target:
        return [], []

    for sibling in target.find_next_siblings():
        if re.match(r"^h[1-4]$", sibling.name or ""):
            break
        if sibling.name == "table":
            headers = [th.get_text(strip=True) for th in sibling.find_all("th")]
            rows = []
            for tr in sibling.find_all("tr")[1:max_rows + 1]:
                cells = [td.get_text(strip=True) for td in tr.find_all("td")]
                if cells:
                    rows.append(cells)
            return headers, rows
    return [], []


def _add_slide_chrome(slide, title: str, customer_name: str = "",
                      logicalis_logo: str = "", subtitle: str = ""):
    """
    Add the standard content-slide chrome:
      - Logicalis logo — top right
      - Red title bar — top left, stopping before the logo
      - Customer name footer — bottom left with blue underline accent
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from lxml import etree

    # Logicalis logo — top right
    if logicalis_logo and os.path.exists(logicalis_logo):
        try:
            slide.shapes.add_picture(logicalis_logo,
                                     _inches(8.3), _inches(0.1),
                                     height=_inches(0.45))
        except Exception:
            pass

    # Red title bar — x=0, full width but leaves visual room for logo
    bar = slide.shapes.add_shape(
        1,
        _inches(0), _inches(0.1),
        _inches(10), _inches(0.65),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(*_RED)
    bar.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = False
    tf.margin_left = _inches(0.15)
    tf.margin_top = _inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size = _pt(18)
    run.font.bold = True
    run.font.color.rgb = _rgb(*_WHITE)
    run.font.name = "Arial"

    # Customer name footer — bottom left with blue underline accent
    if customer_name:
        footer_box = slide.shapes.add_textbox(
            _inches(0.2), _inches(7.1),
            _inches(3.5), _inches(0.3),
        )
        tf_f = footer_box.text_frame
        p_f = tf_f.paragraphs[0]
        run_f = p_f.add_run()
        run_f.text = customer_name
        run_f.font.size = _pt(8)
        run_f.font.color.rgb = _rgb(*_BLUE)
        run_f.font.bold = True
        run_f.font.name = "Arial"

        # Blue bottom border under the footer text
        pPr = p_f._p.get_or_add_pPr()
        pBdr = etree.SubElement(pPr, qn("a:buNone"))  # clear bullet
        lnSpc = etree.SubElement(p_f._p, qn("a:endParaRPr"))


# Keep the old name as an alias used in the generate function
def _add_slide_title(slide, title: str, subtitle: str = ""):
    """Legacy wrapper — use _add_slide_chrome for new slides."""
    _add_slide_chrome(slide, title)


def _add_text_box(slide, text: str, left, top, width, height,
                  font_size: int = 11, bold: bool = False,
                  color: tuple = _DARK, wrap: bool = True):
    from pptx.enum.text import PP_ALIGN
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = wrap
    p = txBox.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = _pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(*color)
    run.font.name = "Arial"
    return txBox


def _add_table_to_slide(slide, headers: list[str], rows: list[list[str]],
                         left, top, width, height):
    """Add a styled table to a slide."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    n_rows = len(rows) + 1  # +1 for header
    n_cols = len(headers) if headers else (len(rows[0]) if rows else 1)
    if n_rows < 2 or n_cols < 1:
        return

    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    # Style header row
    for ci, header in enumerate(headers[:n_cols]):
        cell = table.cell(0, ci)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(*_BLUE)
        p = cell.text_frame.paragraphs[0]
        run = p.runs[0] if p.runs else p.add_run()
        run.font.bold = True
        run.font.color.rgb = _rgb(*_WHITE)
        run.font.size = _pt(9)
        run.font.name = "Arial"

    # Fill data rows
    for ri, row in enumerate(rows):
        bg = (0xF7, 0xF9, 0xFF) if ri % 2 == 0 else _WHITE
        for ci, cell_text in enumerate(row[:n_cols]):
            cell = table.cell(ri + 1, ci)
            cell.text = str(cell_text)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(*bg)
            p = cell.text_frame.paragraphs[0]
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = _pt(8)
            run.font.name = "Arial"
            run.font.color.rgb = _rgb(*_DARK)


def _embed_chart(slide, chart_bytes: bytes, left, top, width, height=None):
    """Embed a chart PNG into a slide."""
    buf = BytesIO(chart_bytes)
    pic = slide.shapes.add_picture(buf, left, top, width=width, height=height)
    return pic


def _embed_logo(slide, logo_path: str, left, top, height):
    """Add a logo image to a slide."""
    if not logo_path or not os.path.exists(logo_path):
        return
    try:
        slide.shapes.add_picture(logo_path, left, top, height=height)
    except Exception as e:
        logger.warning("Could not embed logo %s: %s", logo_path, e)


def _stat_box(slide, label: str, value: str, left, top, width=_inches(2.2), height=_inches(1.2)):
    """Render a metric callout box."""
    from pptx.util import Inches
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = _rgb(0xDB, 0xEA, 0xFE)
    box.line.color.rgb = _rgb(*_BLUE)
    box.line.width = _pt(1)

    tf = box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    run1 = p1.add_run()
    run1.text = value
    run1.font.size = _pt(24)
    run1.font.bold = True
    run1.font.color.rgb = _rgb(*_BLUE)
    run1.font.name = "Arial"

    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = _pt(9)
    run2.font.color.rgb = _rgb(*_GREY)
    run2.font.name = "Arial"


def generate_pptx(markdown_content: str, customer_name: str, report_date: str,
                  logo_path: str | None = None, charts: dict | None = None) -> bytes:
    """
    Generate an executive summary PowerPoint deck.

    Args:
        markdown_content: Full report markdown string
        customer_name: Customer display name
        report_date: "YYYY-MM-DD to YYYY-MM-DD" string
        logo_path: Absolute path to customer logo (optional)
        charts: Dict of chart_name -> PNG bytes

    Returns:
        PPTX file as bytes
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    charts = charts or {}
    report_month = _get_report_month(report_date)

    # Convert markdown to HTML for parsing
    html = md_lib.markdown(markdown_content, extensions=["tables", "fenced_code", "nl2br"])
    soup = BeautifulSoup(html, "html.parser")

    prs = Presentation()
    prs.slide_width = _inches(10)
    prs.slide_height = _inches(7.5)

    blank_layout = prs.slide_layouts[6]  # Blank layout

    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    default_logo   = os.path.join(base_dir, "static", "Logo.webp")
    samples_logo   = os.path.join(base_dir, "samples", "Logicalis_logo_2.png")
    content_logo   = samples_logo if os.path.exists(samples_logo) else default_logo

    def _content_slide(title: str) -> object:
        """Add a blank content slide with standard chrome (logo, red title, footer)."""
        s = prs.slides.add_slide(blank_layout)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = _rgb(*_LIGHT_BG)
        _add_slide_chrome(s, title, customer_name=customer_name,
                          logicalis_logo=content_logo)
        return s

    # ── Slide 1: Cover ────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(*_BLUE)

    # Logicalis logo — top left (white version on blue)
    if os.path.exists(default_logo):
        try:
            slide.shapes.add_picture(default_logo,
                                     _inches(0.35), _inches(0.3),
                                     height=_inches(0.55))
        except Exception:
            pass

    # Red title box 1 — report type
    box1 = slide.shapes.add_shape(
        1, _inches(0.35), _inches(2.1), _inches(6.0), _inches(0.75))
    box1.fill.solid()
    box1.fill.fore_color.rgb = _rgb(*_RED)
    box1.line.fill.background()
    tf1 = box1.text_frame
    tf1.margin_left = _inches(0.1)
    tf1.margin_top  = _inches(0.06)
    p1 = tf1.paragraphs[0]
    r1 = p1.add_run()
    r1.text = "Monthly Security Operations Report"
    r1.font.size = _pt(18)
    r1.font.bold = True
    r1.font.color.rgb = _rgb(*_WHITE)
    r1.font.name = "Arial"

    # Red title box 2 — customer name
    box2 = slide.shapes.add_shape(
        1, _inches(0.35), _inches(3.0), _inches(5.0), _inches(0.75))
    box2.fill.solid()
    box2.fill.fore_color.rgb = _rgb(*_RED)
    box2.line.fill.background()
    tf2 = box2.text_frame
    tf2.margin_left = _inches(0.1)
    tf2.margin_top  = _inches(0.06)
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = customer_name
    r2.font.size = _pt(18)
    r2.font.bold = True
    r2.font.color.rgb = _rgb(*_WHITE)
    r2.font.name = "Arial"

    # Customer logo — below title boxes
    if logo_path and os.path.exists(logo_path):
        try:
            slide.shapes.add_picture(logo_path,
                                     _inches(0.35), _inches(4.0),
                                     height=_inches(1.2))
        except Exception:
            pass

    # Footer bar — bottom left
    footer_bar = slide.shapes.add_shape(
        1, _inches(0), _inches(7.1), _inches(3.5), _inches(0.4))
    footer_bar.fill.solid()
    footer_bar.fill.fore_color.rgb = _rgb(0x1A, 0x4A, 0xBB)
    footer_bar.line.fill.background()
    tf_fb = footer_bar.text_frame
    tf_fb.margin_left = _inches(0.15)
    tf_fb.margin_top  = _inches(0.06)
    p_fb = tf_fb.paragraphs[0]
    r_fb = p_fb.add_run()
    r_fb.text = customer_name
    r_fb.font.size = _pt(9)
    r_fb.font.color.rgb = _rgb(*_WHITE)
    r_fb.font.name = "Arial"

    # ── Slide 2: Executive Summary ────────────────────────────────────────────
    slide = _content_slide("Executive Summary")
    exec_text = _extract_section_text(soup, "Executive Summary", max_chars=700)
    if not exec_text:
        exec_text = _extract_section_text(soup, "Introduction", max_chars=700)
    _add_text_box(slide, exec_text or "See full report for details.",
                  _inches(0.3), _inches(1.0), _inches(9.4), _inches(5.8),
                  font_size=10, color=_DARK)

    # ── Slide 3: Incident Severity ────────────────────────────────────────────
    slide = _content_slide("Incident Severity")
    if charts.get("severity"):
        _embed_chart(slide, charts["severity"],
                     _inches(0.5), _inches(1.0), _inches(9))
    else:
        _add_text_box(slide, "Severity chart not available.",
                      _inches(1), _inches(3), _inches(8), _inches(1),
                      font_size=12, color=_GREY)

    # ── Slide 4: Incident Trend ───────────────────────────────────────────────
    slide = _content_slide("12-Month Incident Trend")
    if charts.get("monthly_trend"):
        _embed_chart(slide, charts["monthly_trend"],
                     _inches(0.5), _inches(1.0), _inches(9))
    else:
        _add_text_box(slide, "Trend chart not available.",
                      _inches(1), _inches(3), _inches(8), _inches(1),
                      font_size=12, color=_GREY)

    # ── Slide 5: Incident Resolution ─────────────────────────────────────────
    slide = _content_slide("Incident Resolution Status")
    if charts.get("resolution"):
        _embed_chart(slide, charts["resolution"],
                     _inches(0.5), _inches(1.0), _inches(9))
    else:
        _add_text_box(slide, "Resolution chart not available.",
                      _inches(1), _inches(3), _inches(8), _inches(1),
                      font_size=12, color=_GREY)

    # ── Slide 6: Top Alerts ───────────────────────────────────────────────────
    slide = _content_slide("Top Alerts Triggered")
    if charts.get("top_alerts"):
        _embed_chart(slide, charts["top_alerts"],
                     _inches(0.3), _inches(1.0), _inches(5.5), _inches(5.8))
    else:
        _add_text_box(slide, "Top alerts chart not available.",
                      _inches(0.3), _inches(1.0), _inches(5.5), _inches(1),
                      font_size=11, color=_GREY)
    alert_headers, alert_rows = _extract_table_rows(soup, "Top Alert", max_rows=5)
    if not alert_headers and not alert_rows:
        alert_headers, alert_rows = _extract_table_rows(soup, "1.7", max_rows=5)
    if alert_headers and alert_rows:
        _add_table_to_slide(slide, alert_headers[:3], [r[:3] for r in alert_rows],
                            _inches(5.9), _inches(1.2), _inches(3.9), _inches(4.5))

    # ── Slide 7: Pending Tickets ──────────────────────────────────────────────
    pending_headers, pending_rows = _extract_table_rows(soup, "Pending Ticket", max_rows=8)
    if not pending_headers:
        pending_headers, pending_rows = _extract_table_rows(soup, "1.14", max_rows=8)
    slide = _content_slide("Pending Tickets")
    if pending_headers and pending_rows:
        _add_table_to_slide(slide, pending_headers[:5], [r[:5] for r in pending_rows],
                            _inches(0.3), _inches(1.1), _inches(9.4), _inches(5.7))
    else:
        _add_text_box(slide, "No pending tickets during this period.",
                      _inches(1), _inches(3), _inches(8), _inches(1),
                      font_size=13, color=_GREY)

    # ── Slide 8: SOCRadar Threat Intelligence (if available) ──────────────────
    socradar_text = _extract_section_text(soup, "SOCRadar", max_chars=700)
    if socradar_text and "Data Source Pending" not in socradar_text:
        slide = _content_slide("SOCRadar Threat Intelligence")
        _add_text_box(slide, socradar_text,
                      _inches(0.3), _inches(1.0), _inches(9.4), _inches(5.8),
                      font_size=10, color=_DARK)

    # ── Slide 9: Recommendations ──────────────────────────────────────────────
    rec_headers, rec_rows = _extract_table_rows(soup, "Recommendation", max_rows=6)
    slide = _content_slide("GSOC Recommendations")
    if rec_headers and rec_rows:
        _add_table_to_slide(slide, rec_headers[:5], [r[:5] for r in rec_rows],
                            _inches(0.3), _inches(1.1), _inches(9.4), _inches(5.7))
    else:
        rec_text = _extract_section_text(soup, "Recommendation", max_chars=600)
        _add_text_box(slide, rec_text or "No recommendations data available.",
                      _inches(0.3), _inches(1.0), _inches(9.4), _inches(5.8),
                      font_size=10, color=_DARK)

    # ── Slide 10: Monitoring Scope + Confidentiality ──────────────────────────
    slide = _content_slide("Monitoring Scope & Confidentiality")
    scope_text = _extract_section_text(soup, "Monitoring Scope", max_chars=400)
    _add_text_box(slide, scope_text or "Refer to the full report for monitoring scope details.",
                  _inches(0.3), _inches(1.0), _inches(9.4), _inches(2.5),
                  font_size=10, color=_DARK)
    _add_text_box(slide,
                  "The contents of this document are confidential and proprietary to Logicalis. "
                  "Do not disclose to any third party without written consent.",
                  _inches(0.3), _inches(4.5), _inches(9.4), _inches(1.2),
                  font_size=9, color=_GREY)

    # Serialize to bytes
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
